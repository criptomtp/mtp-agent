"""ContentAgent — генерація HTML презентації, PPTX та email тексту."""

import os
import re
import logging
import requests
from typing import Dict, Any, Optional
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .research_agent import Lead

logger = logging.getLogger(__name__)

HARDCODED_TARIFFS_TABLE = [
    ("Прийом товару", "2 грн / одиниця"),
    ("Зберігання (паллет)", "800 грн / місяць"),
    ("Зберігання (коробка)", "80 грн / місяць"),
    ("Комплектація B2C (1-3 од.)", "22 грн / замовлення"),
    ("Комплектація B2B", "45 грн / замовлення"),
    ("Пакування", "8 грн / замовлення"),
    ("Відправка Новою Поштою", "за тарифом НП"),
]


def _build_tariffs_rows(tariffs=None):
    """Build tariff rows from DB tariffs or fallback to hardcoded."""
    if not tariffs:
        return HARDCODED_TARIFFS_TABLE

    rows = []
    for t in tariffs:
        name = t["service_name"]
        price = t["price"]
        unit = t.get("unit", "")
        note = t.get("note", "")
        if price > 0:
            tariff_str = f"{int(price)} грн / {unit}" if unit else f"{int(price)} грн"
        elif note:
            tariff_str = note
        else:
            continue
        rows.append((name, tariff_str))
    return rows if rows else HARDCODED_TARIFFS_TABLE


def _escape_html(text: str) -> str:
    """Escape HTML special characters."""
    return (text or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;")


def _get_text_color_for_bg(hex_color: str) -> str:
    """Return '#fff' or '#222' based on perceived brightness of background color."""
    hex_color = hex_color.strip().lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join(c * 2 for c in hex_color)
    try:
        r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    except (ValueError, IndexError):
        return "#fff"
    brightness = (r * 299 + g * 587 + b * 114) / 1000
    return "#222" if brightness > 160 else "#fff"


class ContentAgent:
    """Генерує HTML презентацію та email текст."""

    def __init__(self, api_keys: Optional[Dict[str, str]] = None):
        self._api_keys = api_keys or {}

    def generate(self, lead: Lead, analysis: Dict[str, Any], output_dir: str,
                 tariffs=None, brand_style: Optional[Dict[str, str]] = None, niche: str = "") -> Dict[str, str]:
        """Генерує HTML, PPTX та email.txt. Повертає шляхи до файлів."""
        os.makedirs(output_dir, exist_ok=True)

        html_path = os.path.join(output_dir, "proposal.html")
        email_path = os.path.join(output_dir, "email.txt")
        safe_name = re.sub(r"[^\w\s-]", "", lead.name).strip().replace(" ", "_")[:50]
        pptx_path = os.path.join(output_dir, f"{safe_name}_proposal.pptx")

        html = self._generate_html_proposal(lead, analysis, tariffs, brand_style=brand_style)
        with open(html_path, "w", encoding="utf-8") as f:
            f.write(html)
        logger.info(f"HTML presentation saved: {html_path}")

        # PPTX generation
        try:
            prs = self._generate_pptx_proposal(lead, analysis, tariffs)
            prs.save(pptx_path)
            logger.info(f"PPTX presentation saved: {pptx_path}")
        except Exception as e:
            logger.error(f"PPTX generation failed: {e}", exc_info=True)
            pptx_path = ""

        self._generate_email(lead, analysis, email_path)

        # Web proposal — Gemini generates unique HTML, fallback to template
        web_proposal = None
        try:
            web_proposal = self.create_web_proposal(lead, analysis, brand_style=brand_style,
                                                     tariffs=tariffs, niche=niche)
        except Exception as e:
            logger.error(f"Web proposal creation failed: {e}", exc_info=True)

        if not web_proposal:
            try:
                web_proposal = self._create_web_proposal_fallback(lead, analysis, brand_style=brand_style, tariffs=tariffs)
            except Exception as e:
                logger.error(f"Web proposal fallback also failed: {e}", exc_info=True)

        # Overwrite local proposal.html with Gemini HTML if available
        if web_proposal and web_proposal.get("html_content"):
            with open(html_path, "w", encoding="utf-8") as f:
                f.write(web_proposal["html_content"])
            logger.info(f"[ContentAgent] Overwrote proposal.html with Gemini HTML")

        result = {"html": html_path, "email": email_path, "pptx": pptx_path}
        if web_proposal:
            result["web_url"] = web_proposal["url"]
            result["web_proposal"] = web_proposal
        return result

    def _validate_html(self, html: str, brand_primary: str = "") -> tuple:
        """Validate generated HTML for quality issues. Returns (is_valid, issues)."""
        issues = []
        if "viewport" not in html:
            issues.append("Missing viewport meta tag")
        # Check for fixed widths > 600px (exclude max-width and media queries)
        fixed_widths = re.findall(r"(?<!max-)(?<!min-)width:\s*(\d+)px", html)
        wide = [w for w in fixed_widths if int(w) > 600]
        if wide:
            issues.append(f"Fixed widths > 600px: {wide[:5]}")
        if "@media" not in html:
            issues.append("No media queries — not responsive")
        if brand_primary and brand_primary.lower() not in html.lower():
            issues.append(f"brand_primary {brand_primary} not found in HTML")
        if "box-sizing" not in html:
            issues.append("Missing box-sizing: border-box")
        return len(issues) == 0, issues

    def create_web_proposal(self, lead: "Lead", analysis: Dict[str, Any],
                            brand_style: Optional[Dict[str, str]] = None,
                            tariffs=None, niche: str = "") -> Optional[Dict[str, str]]:
        """Generate unique HTML proposal via Gemini AI and upload to Supabase Storage.

        Returns {'slug': ..., 'url': ..., 'proposal_id': ...} or None on error.
        """
        import uuid
        import hashlib

        api_key = self._api_keys.get("GEMINI_API_KEY") or os.getenv("GEMINI_API_KEY")
        if not api_key:
            logger.error("[ContentAgent] GEMINI_API_KEY not set, cannot generate web proposal")
            return None

        # Build analysis text
        pain_points = analysis.get("pain_points", [])
        pain_titles = []
        for p in pain_points:
            if isinstance(p, dict):
                pain_titles.append(p.get("title", ""))
            else:
                pain_titles.append(str(p))

        benefits = analysis.get("key_benefits", [])
        benefit_titles = []
        for b in benefits:
            if isinstance(b, dict):
                benefit_titles.append(b.get("benefit", ""))
            else:
                benefit_titles.append(str(b))

        score = analysis.get("score", 0)
        grade = analysis.get("grade", "?")
        analysis_text = (
            f"Потенціал: {score}/10 ({grade}). "
            f"Болі: {', '.join(pain_titles)}. "
            f"Переваги MTP для клієнта: {', '.join(benefit_titles)}"
        )

        from .knowledge_base import MTP_COMPANY, MTP_CLIENTS, MTP_TARIFFS, get_tariffs_prompt_text, get_clients_text

        # Brand style
        bs = brand_style or {}
        brand_primary = bs.get("primary_color", "#1A365D")
        brand_secondary = bs.get("secondary_color", "#E53E3E")
        brand_font = bs.get("font_family", "Inter")

        # Generate slug
        slug = hashlib.md5(f"{lead.name}-{datetime.now().isoformat()}".encode()).hexdigest()[:12]

        client_name = lead.name
        website = getattr(lead, "website", "") or ""
        city = getattr(lead, "city", "") or ""
        niche_text = niche or "e-commerce"
        api_base = os.getenv("MTP_API_URL", "https://mtp-agent-production.up.railway.app")
        calendly_url = MTP_COMPANY["calendly"]
        tariffs_prompt = get_tariffs_prompt_text()
        clients_text = get_clients_text()
        clients_links = " ".join([f'<a href="{c["url"]}" target="_blank">{c["name"]}</a>' for c in MTP_CLIENTS])

        prompt = f"""<BRAND_COLORS>
primary_color: {brand_primary}
secondary_color: {brand_secondary}
font: {brand_font}
</BRAND_COLORS>

ІНСТРУКЦІЯ: hero фон = {brand_primary}, navbar фон = {brand_primary}, заголовки h2 колір = {brand_primary}.
НЕ використовуй #1A365D якщо brand_primary != #1A365D. НЕ вигадуй свої кольори.

Ти — топовий UX/UI дизайнер і копірайтер. Створи повну HTML-сторінку комерційної пропозиції від MTP Fulfillment для клієнта.

КЛІЄНТ:
- Назва: {client_name}
- Ніша: {niche_text}
- Сайт: {website}
- Місто: {city}
- Аналіз (болі, потенціал): {analysis_text}
- Фірмовий стиль клієнта: primary={brand_primary}, secondary={brand_secondary}, font={brand_font}

MTP FULFILLMENT (продавець):
- Компанія: {MTP_COMPANY['name']}, {MTP_COMPANY['location']}
- Колір бренду: #E53E3E (червоний), білий
- {MTP_COMPANY['years']} років на ринку, {MTP_COMPANY['shipments_per_month']} відправок/міс, {MTP_COMPANY['warehouses']} склади
- Контакти: {MTP_COMPANY['email']}, {MTP_COMPANY['phone']}, {MTP_COMPANY['website']}
- Calendly: {calendly_url}
- Клієнти-референси: {clients_text}

{tariffs_prompt}

ЗАБОРОНЕНО: НЕ вигадуй тарифи! Використовуй ВИКЛЮЧНО значення з тарифної сітки вище.
Якщо не знаєш тариф — пиши "за запитом" або не пиши взагалі.
НЕ змішуй категорії: зберігання ≠ відвантаження ≠ приймання.

ЗАВДАННЯ:
Створи УНІКАЛЬНУ HTML сторінку яка:
1. Використовує фірмові кольори КЛІЄНТА як основну палітру ({brand_primary} як домінантний)
2. Додає MTP червоний (#E53E3E) як акцентний колір на CTA кнопках і важливих елементах
3. Має УНІКАЛЬНУ структуру і layout — НЕ стандартний шаблон. Придумай щось оригінальне для цього клієнта
4. Пише УНІКАЛЬНИЙ текст під цю нішу і цього клієнта — не загальні фрази
5. Включає секції (але в унікальному порядку і стилі): hero з назвою клієнта, болі бізнесу, рішення від MTP, тарифи (ТОЧНО за структурою вище), соціальний доказ ({clients_links}), CTA
6. Адаптивна (mobile-friendly)
7. Без зовнішніх залежностей (тільки Google Fonts дозволені)
8. Tracking: при завантаженні fetch('{api_base}/api/proposals/track', {{method:'POST', headers:{{'Content-Type':'application/json'}}, body: JSON.stringify({{slug:'{slug}', event:'open'}})}})\

ОБОВ'ЯЗКОВІ ВИМОГИ ДО ВЕРСТКИ:
- Sticky navbar (position:fixed, top:0, z-index:100, height ~60-70px) з: логотип "MTP Fulfillment" → <a href="{MTP_COMPANY['website']}" target="_blank">, телефон → <a href="tel:{MTP_COMPANY['phone_raw']}">{MTP_COMPANY['phone']}</a>, кнопка "Записатись" → <a href="{calendly_url}" target="_blank">
- Hero секція ОБОВ'ЯЗКОВО має padding-top: мінімум 100px (щоб текст не ховався під sticky navbar)
- Всі body секції мають padding-top достатній щоб не перекриватись navbar
- Клієнти-референси — ТІЛЬКИ текстові бейджі (НЕ використовуй img, НЕ placeholder.com, НЕ via.placeholder): {clients_links}
- Footer — всі контакти клікабельні: <a href="mailto:{MTP_COMPANY['email']}">, <a href="tel:{MTP_COMPANY['phone_raw']}">, <a href="{MTP_COMPANY['website']}" target="_blank">. Копірайт: © 2026 MTP Fulfillment
- Перевір: весь текст читабельний, немає overflow, кнопки не перекривають текст, контраст достатній

АБСОЛЮТНА ВИМОГА ДО КОЛЬОРІВ (ПОРУШЕННЯ = БРАК):
- CSS змінна --brand-primary ПОВИННА бути: {brand_primary}
- Hero секція, navbar, всі заголовки h2 — фон або колір тексту {brand_primary}
- НЕ ВИКОРИСТОВУЙ #1A365D, #1a1a2e, #2c3e50 або будь-який інший колір замість {brand_primary}
- Якщо brand_primary = #1A365D — це дефолтний колір, використай нейтральний (#0f172a)
- MTP червоний #E53E3E — ТІЛЬКИ для CTA кнопок та акцентів, НЕ для фону секцій
- ЗАБОРОНЕНО: ігнорувати brand_primary або замінювати його на свій вибір
- Перевір перед відповіддю: чи є {brand_primary} в CSS root і чи він реально використовується

КРИТИЧНО — АДАПТИВНІСТЬ І ЧИТАБЕЛЬНІСТЬ:
- * {{ box-sizing: border-box; }}
- img {{ max-width: 100%; height: auto; }}
- Кнопки: white-space: nowrap; max-width: 100%; overflow: hidden; text-overflow: ellipsis
- Navbar кнопки на мобільному: font-size: 14px; padding: 8px 16px
- Таблиці тарифів: width: 100%; table-layout: fixed; word-wrap: break-word
- Текст в картках: overflow: hidden; word-break: break-word
- НЕ використовуй фіксовані ширини в px для кнопок і карток — тільки %, vw, max-width
- Медіа запити ОБОВ'ЯЗКОВІ: @media (max-width: 768px) і @media (max-width: 480px)
- Перевір: чи всі кнопки поміщаються на екрані 375px шириною?

ВАЖЛИВО: Поверни ТІЛЬКИ валідний HTML від <!DOCTYPE html> до </html>. Без markdown, без пояснень. Рік у копірайті — 2026."""

        try:
            import google.generativeai as genai
            genai.configure(api_key=api_key)

            # Read model from pipeline settings, fallback chain if it fails
            configured_model = "gemini-2.0-flash"
            try:
                import json as _json
                settings_path = os.path.join(os.path.dirname(__file__), "..", "config", "pipeline_settings.json")
                with open(settings_path, "r", encoding="utf-8") as _f:
                    _settings = _json.load(_f)
                configured_model = _settings.get("agents", {}).get("content", {}).get("model", "gemini-2.0-flash")
                logger.info(f"[ContentAgent] Using model from settings: {configured_model}")
            except Exception:
                logger.info(f"[ContentAgent] No settings file, using default: {configured_model}")

            # Try configured model first, then fallbacks
            models_to_try = [configured_model]
            for fallback in ["gemini-2.5-flash", "gemini-2.0-flash", "gemini-1.5-flash"]:
                if fallback not in models_to_try:
                    models_to_try.append(fallback)

            html = None
            for model_name in models_to_try:
                try:
                    logger.info(f"[ContentAgent] Trying {model_name} for web proposal...")
                    model = genai.GenerativeModel(model_name)
                    response = model.generate_content(prompt)
                    if response.text:
                        html = response.text.strip()
                        logger.info(f"[ContentAgent] {model_name} returned {len(html)} bytes")
                        break
                    else:
                        logger.warning(f"[ContentAgent] {model_name} returned empty response")
                except Exception as model_err:
                    logger.warning(f"[ContentAgent] {model_name} failed: {model_err}")
                    continue

            if not html:
                logger.error("[ContentAgent] All Gemini models failed for web proposal")
                return None

            # Strip markdown fences if present
            if html.startswith("```"):
                html = html.split("\n", 1)[1].rsplit("```", 1)[0].strip()

            # Validate HTML
            if not html.startswith("<!DOCTYPE html>") and not html.startswith("<!doctype html>"):
                # Try to find DOCTYPE in response
                idx = html.lower().find("<!doctype html>")
                if idx >= 0:
                    html = html[idx:]
                else:
                    logger.error("[ContentAgent] Gemini response is not valid HTML")
                    return None

            logger.info(f"[ContentAgent] Gemini generated {len(html)} bytes of unique HTML for {client_name}")

            # Validate HTML quality
            valid, issues = self._validate_html(html, brand_primary)
            if not valid:
                logger.warning(f"[ContentAgent] HTML issues for {client_name}: {'; '.join(issues)}")

            # Force brand color if Gemini ignored it
            if brand_primary.lower() != "#1a365d" and brand_primary.lower() not in html.lower():
                logger.warning(f"[ContentAgent] Gemini ignored brand color {brand_primary} for {client_name}, force-replacing")
                # Replace common default colors Gemini might use instead
                for default_color in ["#1A365D", "#1a365d", "#2c3e50", "#1a1a2e", "#0f172a"]:
                    if default_color in html:
                        html = html.replace(default_color, brand_primary)
                        logger.info(f"[ContentAgent] Replaced {default_color} → {brand_primary}")
                        break

            # Inject CSS overrides for consistent rendering
            css_override = """<style>
* { font-family: 'Inter', 'Segoe UI', Arial, Helvetica, sans-serif !important; }
table { table-layout: fixed !important; width: 100% !important; }
td, th { overflow-wrap: break-word !important; word-break: break-word !important; }
img { max-width: 100% !important; height: auto !important; }
</style>"""
            html = html.replace("</head>", css_override + "\n</head>", 1)

        except Exception as e:
            logger.error(f"[ContentAgent] Gemini HTML generation failed: {e}", exc_info=True)
            return None

        # Save HTML directly to proposals table (Supabase Storage has encoding issues)
        try:
            from backend.services.database import get_supabase_admin

            db = get_supabase_admin()
            proposal_data = {
                "slug": slug,
                "client_name": client_name,
                "html_content": html,
                "client_data": {
                    "niche": niche_text,
                    "website": website,
                    "city": city,
                    "score": score,
                    "grade": grade,
                    "brand_style": bs,
                },
            }
            try:
                result = db.table("proposals").insert(proposal_data).execute()
                if result.data:
                    logger.info(f"[ContentAgent] Proposal saved to DB: slug={slug}, id={result.data[0].get('id')}")
                else:
                    logger.error(f"[ContentAgent] Proposal insert returned no data for slug={slug}")
            except Exception as e:
                logger.error(f"[ContentAgent] Proposal insert failed: {e}", exc_info=True)
                # Retry without html_content if column doesn't exist yet
                try:
                    proposal_data.pop("html_content", None)
                    proposal_data["client_data"]["html_content"] = html
                    db.table("proposals").insert(proposal_data).execute()
                    logger.info(f"[ContentAgent] Proposal saved with html in client_data: slug={slug}")
                except Exception as e2:
                    logger.error(f"[ContentAgent] Proposal insert retry failed: {e2}", exc_info=True)

            proposal_url = f"{api_base}/api/proposals/{slug}"
            logger.info(f"[ContentAgent] Web proposal ready: {proposal_url}")
            return {"slug": slug, "url": proposal_url, "proposal_id": slug, "html_content": html}

        except Exception as e:
            logger.error(f"[ContentAgent] Proposal save failed: {e}", exc_info=True)
            return None

    def _create_web_proposal_fallback(self, lead: "Lead", analysis: Dict[str, Any],
                                      brand_style: Optional[Dict[str, str]] = None,
                                      tariffs=None) -> Optional[Dict[str, str]]:
        """Fallback: save template HTML to proposals table if Gemini fails."""
        import hashlib

        try:
            from backend.services.database import get_supabase_admin

            html = self._generate_html_proposal(lead, analysis, tariffs, brand_style=brand_style)
            slug = hashlib.md5(f"{lead.name}-fallback-{datetime.now().isoformat()}".encode()).hexdigest()[:12]
            api_base = os.getenv("MTP_API_URL", "https://mtp-agent-production.up.railway.app")

            db = get_supabase_admin()
            proposal_data = {
                "slug": slug,
                "client_name": lead.name,
                "html_content": html,
            }
            try:
                result = db.table("proposals").insert(proposal_data).execute()
                if result.data:
                    logger.info(f"[ContentAgent] Fallback proposal saved: slug={slug}")
            except Exception:
                # Retry with html in client_data if html_content column missing
                try:
                    db.table("proposals").insert({
                        "slug": slug,
                        "client_name": lead.name,
                        "client_data": {"html_content": html},
                    }).execute()
                    logger.info(f"[ContentAgent] Fallback proposal saved in client_data: slug={slug}")
                except Exception as e:
                    logger.error(f"[ContentAgent] Fallback proposal insert failed: {e}", exc_info=True)
                    return None

            proposal_url = f"{api_base}/api/proposals/{slug}"
            logger.info(f"[ContentAgent] Fallback proposal ready: {proposal_url}")
            return {"slug": slug, "url": proposal_url, "proposal_id": slug, "html_content": html}

        except Exception as e:
            logger.error(f"[ContentAgent] Fallback save failed: {e}", exc_info=True)
            return None

    # ── PPTX helpers ──

    @staticmethod
    def _pptx_add_bg(slide, prs, color):
        """Add a full-slide background rectangle."""
        from pptx.util import Emu
        from pptx.enum.shapes import MSO_SHAPE
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()

    @staticmethod
    def _pptx_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                       color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.LEFT, font_name="Calibri"):
        """Add a text box with styled text."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        return tf

    def _generate_pptx_proposal(self, lead: Lead, analysis: Dict[str, Any], tariffs=None):
        """Generate 6-slide dark-themed PPTX presentation."""
        BG = RGBColor(0x0F, 0x20, 0x44)
        BG_CARD = RGBColor(0x16, 0x2D, 0x5A)
        ACCENT = RGBColor(0x00, 0xB4, 0xD8)
        MINT = RGBColor(0x06, 0xD6, 0xA0)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        LIGHT = RGBColor(0xBB, 0xBB, 0xCC)
        ORANGE = RGBColor(0xFF, 0x6B, 0x35)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        W = prs.slide_width
        H = prs.slide_height
        MARGIN = Inches(0.8)
        CONTENT_W = W - 2 * MARGIN

        date_str = datetime.now().strftime("%d.%m.%Y")
        hook = analysis.get("hook", f"{lead.name}: час масштабуватись")
        client_insight = analysis.get("client_insight", "")
        zoom_cta = analysis.get("zoom_cta", "Запишіться на безкоштовну Zoom-консультацію!")

        # ── SLIDE 1: Hero ──
        slide = prs.slides.add_slide(blank_layout)
        self._pptx_add_bg(slide, prs, BG)
        self._pptx_text_box(slide, MARGIN, Inches(1.0), CONTENT_W, Inches(0.6),
                            "MTP FULFILLMENT", font_size=14, color=ACCENT, bold=True)
        self._pptx_text_box(slide, MARGIN, Inches(2.0), CONTENT_W, Inches(1.5),
                            lead.name, font_size=48, bold=True, color=WHITE)
        subtitle = f"Комерційна пропозиція"
        if lead.city:
            subtitle += f"  •  {lead.city}"
        self._pptx_text_box(slide, MARGIN, Inches(3.8), CONTENT_W, Inches(0.6),
                            subtitle, font_size=20, color=LIGHT)
        self._pptx_text_box(slide, MARGIN, Inches(6.2), CONTENT_W, Inches(0.5),
                            date_str, font_size=14, color=LIGHT, alignment=PP_ALIGN.RIGHT)
        from pptx.enum.shapes import MSO_SHAPE

        # ── SLIDE 2: Insight ──
        slide = prs.slides.add_slide(blank_layout)
        self._pptx_add_bg(slide, prs, BG)
        self._pptx_text_box(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.4),
                            "ПРО НАС", font_size=12, color=ACCENT, bold=True)
        # Stats row
        stats = [("7+", "років на ринку"), ("60K+", "відправок / міс"), ("2", "склади під Києвом"), ("30с", "середній час обробки")]
        stat_w = Inches(2.8)
        for i, (num, label) in enumerate(stats):
            x = MARGIN + Emu(int(i * stat_w))
            self._pptx_text_box(slide, x, Inches(1.2), stat_w, Inches(0.8),
                                num, font_size=44, bold=True, color=MINT)
            self._pptx_text_box(slide, x, Inches(2.0), stat_w, Inches(0.5),
                                label, font_size=13, color=LIGHT)
        # Insight text
        if client_insight:
            self._pptx_text_box(slide, MARGIN, Inches(3.2), CONTENT_W, Inches(0.4),
                                "МИ РОЗУМІЄМО ВАШУ СПЕЦИФІКУ", font_size=12, color=ACCENT, bold=True)
            self._pptx_text_box(slide, MARGIN, Inches(3.8), CONTENT_W, Inches(2.5),
                                client_insight, font_size=18, color=WHITE)

        # ── SLIDE 3: Pain Points ──
        slide = prs.slides.add_slide(blank_layout)
        self._pptx_add_bg(slide, prs, BG)
        self._pptx_text_box(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.4),
                            "ДЕ ВТРАЧАЄТЬСЯ ЕФЕКТИВНІСТЬ", font_size=12, color=ACCENT, bold=True)
        pain_points = analysis.get("pain_points", [])
        card_w = Inches(3.6)
        card_h = Inches(3.8)
        gap = Inches(0.4)
        for i, p in enumerate(pain_points[:3]):
            x = MARGIN + Emu(int(i * (card_w + gap)))
            # Card background
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), card_w, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = BG_CARD
            card.line.fill.background()
            # Number
            self._pptx_text_box(slide, x + Inches(0.3), Inches(1.6), Inches(1), Inches(0.7),
                                f"0{i+1}", font_size=36, bold=True, color=ACCENT)
            # Title & description
            if isinstance(p, dict):
                title = p.get("title", "")
                desc = p.get("description", "")
            else:
                title = str(p)
                desc = ""
            if len(desc) > 150:
                desc = desc[:147] + "..."
            self._pptx_text_box(slide, x + Inches(0.3), Inches(2.5), card_w - Inches(0.6), Inches(0.6),
                                title, font_size=16, bold=True, color=WHITE)
            if desc:
                self._pptx_text_box(slide, x + Inches(0.3), Inches(3.2), card_w - Inches(0.6), Inches(1.5),
                                    desc, font_size=13, color=LIGHT)

        # ── SLIDE 4: Benefits ──
        slide = prs.slides.add_slide(blank_layout)
        self._pptx_add_bg(slide, prs, BG)
        self._pptx_text_box(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.4),
                            "КЛЮЧОВІ ПЕРЕВАГИ", font_size=12, color=ACCENT, bold=True)
        benefits = analysis.get("key_benefits", [])
        b_w = Inches(5.5)
        b_h = Inches(2.2)
        b_gap_x = Inches(0.5)
        b_gap_y = Inches(0.4)
        for i, kb in enumerate(benefits[:4]):
            col = i % 2
            row = i // 2
            x = MARGIN + Emu(int(col * (b_w + b_gap_x)))
            y = Inches(1.3) + Emu(int(row * (b_h + b_gap_y)))
            # Card
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, b_w, b_h)
            card.fill.solid()
            card.fill.fore_color.rgb = BG_CARD
            card.line.fill.background()
            # Accent stripe
            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Pt(5), b_h)
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = MINT
            stripe.line.fill.background()
            if isinstance(kb, dict):
                benefit = kb.get("benefit", "")
                proof = kb.get("proof", "")
            else:
                benefit = str(kb)
                proof = ""
            self._pptx_text_box(slide, x + Inches(0.4), y + Inches(0.3), b_w - Inches(0.7), Inches(0.6),
                                benefit, font_size=16, bold=True, color=WHITE)
            if proof:
                self._pptx_text_box(slide, x + Inches(0.4), y + Inches(1.0), b_w - Inches(0.7), Inches(1.0),
                                    proof, font_size=13, color=LIGHT)

        # ── SLIDE 5: Pricing ──
        slide = prs.slides.add_slide(blank_layout)
        self._pptx_add_bg(slide, prs, BG)
        self._pptx_text_box(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.4),
                            "ТАРИФИ MTP FULFILLMENT", font_size=12, color=ACCENT, bold=True)
        tariff_rows = _build_tariffs_rows(tariffs)
        rows_count = len(tariff_rows) + 1
        cols = 2
        # Scale row height to fit slide — cap at 0.40 per row, shrink if >12 rows
        max_tbl_height = Inches(5.2)  # leave room for header + bottom
        row_h = min(0.40, 5.2 / max(rows_count, 1))
        tbl_h = Inches(row_h * rows_count)
        tbl_w = Inches(6.5)
        tbl = slide.shapes.add_table(rows_count, cols, MARGIN, Inches(1.3), tbl_w, tbl_h).table
        tbl.columns[0].width = Inches(4.2)
        tbl.columns[1].width = Inches(2.3)
        font_sz = Pt(12) if rows_count <= 12 else Pt(10)
        # Header
        for ci, txt in enumerate(["Послуга", "Тариф"]):
            cell = tbl.cell(0, ci)
            cell.text = txt
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.font.name = "Calibri"
        # Rows
        for ri, (name, price) in enumerate(tariff_rows):
            for ci, txt in enumerate([name, price]):
                cell = tbl.cell(ri + 1, ci)
                cell.text = txt
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_CARD if ri % 2 == 0 else BG
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                for p in cell.text_frame.paragraphs:
                    p.font.size = font_sz
                    p.font.color.rgb = WHITE if ci == 0 else MINT
                    p.font.name = "Calibri"
        # Pricing estimate — positioned to the right of the table, compact 3-line format
        pricing = analysis.get("pricing_estimate", {})
        est_x = MARGIN + tbl_w + Inches(0.5)
        est_w = W - est_x - MARGIN
        if pricing and est_w > Inches(2):
            self._pptx_text_box(slide, est_x, Inches(1.3), est_w, Inches(0.4),
                                "ОРІЄНТОВНИЙ КОШТОРИС", font_size=12, color=ACCENT, bold=True)
            y_off = Inches(2.0)
            if isinstance(pricing, dict):
                # Show max 4 items, truncate values to 50 chars
                items = list(pricing.items())[:4]
                for key, val in items:
                    label = key.replace("_", " ").capitalize()
                    val_str = str(val)[:50]
                    line = f"{label}: {val_str}"
                    self._pptx_text_box(slide, est_x, y_off, est_w - Inches(0.3), Inches(0.5),
                                        line, font_size=11, color=MINT, bold=True)
                    y_off += Inches(0.55)
            else:
                est_text = str(pricing)[:150]
                self._pptx_text_box(slide, est_x, y_off, est_w - Inches(0.3), Inches(1),
                                    est_text, font_size=11, color=WHITE)

        # ── SLIDE 6: CTA ──
        slide = prs.slides.add_slide(blank_layout)
        self._pptx_add_bg(slide, prs, BG)
        self._pptx_text_box(slide, MARGIN, Inches(1.2), CONTENT_W, Inches(0.6),
                            "НАСТУПНИЙ КРОК", font_size=14, color=ACCENT, bold=True)
        cta_short = zoom_cta[:180] + ("..." if len(zoom_cta) > 180 else "")
        self._pptx_text_box(slide, MARGIN, Inches(2.0), CONTENT_W, Inches(2.0),
                            cta_short, font_size=22, bold=True, color=WHITE)
        # CTA button shape
        btn_w = Inches(4)
        btn_h = Inches(0.8)
        btn_x = Emu(int((W - btn_w) / 2))
        btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, btn_x, Inches(4.4), btn_w, btn_h)
        btn.fill.solid()
        btn.fill.fore_color.rgb = ACCENT
        btn.line.fill.background()
        tf = btn.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "ZOOM-ДЗВІНОК"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(10)
        # Contacts
        contacts = [
            "mtpgrouppromo@gmail.com  |  +38 (050) 144-46-45",
            "fulfillmentmtp.com.ua  |  @nikolay_mtp",
        ]
        self._pptx_text_box(slide, MARGIN, Inches(5.3), CONTENT_W, Inches(0.5),
                            contacts[0], font_size=14, color=LIGHT, alignment=PP_ALIGN.CENTER)
        self._pptx_text_box(slide, MARGIN, Inches(5.8), CONTENT_W, Inches(0.5),
                            contacts[1], font_size=14, color=LIGHT, alignment=PP_ALIGN.CENTER)
        self._pptx_text_box(slide, MARGIN, Inches(6.5), CONTENT_W, Inches(0.4),
                            "MTP Fulfillment — 7+ років на ринку, 60 000+ відправок на місяць",
                            font_size=11, color=RGBColor(0x66, 0x66, 0x88), alignment=PP_ALIGN.CENTER)

        return prs

    def _generate_html_proposal(self, lead: Lead, analysis: Dict[str, Any], tariffs=None,
                                brand_style=None) -> str:
        """Generate HTML string for the commercial proposal (modern web design)."""
        from .knowledge_base import MTP_COMPANY, MTP_CLIENTS, get_clients_text

        bs = brand_style or {}
        brand_primary = bs.get("primary_color", "#0f172a")
        brand_secondary = bs.get("secondary_color", "#1e40af")
        mtp_accent = "#E53E3E"
        hero_text_color = _get_text_color_for_bg(brand_primary)
        navbar_text_color = hero_text_color

        hook = _escape_html(analysis.get("hook", f"{lead.name}: час масштабуватись"))
        client_insight = _escape_html(analysis.get("client_insight", ""))
        mtp_fit = _escape_html(analysis.get("mtp_fit", ""))
        zoom_cta = _escape_html(analysis.get("zoom_cta", "Запишіться на безкоштовну Zoom-консультацію!"))
        client_name = _escape_html(lead.name)
        calendly_url = MTP_COMPANY.get("calendly", "https://calendly.com/mtp-fulfillment")

        # Pain points HTML
        pains_html = ""
        for p in analysis.get("pain_points", []):
            if isinstance(p, dict):
                title = _escape_html(p.get("title", ""))
                desc = _escape_html(p.get("description", ""))
                pains_html += f'<div class="pain-card"><h3>{title}</h3><p>{desc}</p></div>'
            else:
                pains_html += f'<div class="pain-card"><p>{_escape_html(str(p))}</p></div>'

        # Key benefits HTML
        benefits_html = ""
        for kb in analysis.get("key_benefits", []):
            if isinstance(kb, dict):
                benefit = _escape_html(kb.get("benefit", ""))
                proof = _escape_html(kb.get("proof", ""))
                benefits_html += f'<div class="benefit-card"><h3>{benefit}</h3><p>{proof}</p></div>'
            else:
                benefits_html += f'<div class="benefit-card"><h3>{_escape_html(str(kb))}</h3></div>'

        # Tariffs table rows
        tariff_rows = _build_tariffs_rows(tariffs)
        tariffs_html = ""
        for name, price in tariff_rows:
            tariffs_html += f"<tr><td>{_escape_html(name)}</td><td>{_escape_html(price)}</td></tr>"

        # Pricing estimate
        pricing_html = ""
        pricing = analysis.get("pricing_estimate", {})
        if pricing:
            pricing_html = '<h2>Орієнтовний кошторис</h2><table class="tariff-table">'
            items = list(pricing.items())
            for i, (key, val) in enumerate(items):
                label = key.replace("_", " ").capitalize()
                cls = ' class="total-row"' if i == len(items) - 1 else ""
                pricing_html += f"<tr{cls}><td>{_escape_html(label)}</td><td>{_escape_html(str(val))}</td></tr>"
            pricing_html += "</table>"

        # Clients references
        clients_html = ""
        for c in MTP_CLIENTS:
            clients_html += f'<a href="{c.get("url","#")}" target="_blank" class="client-logo">{_escape_html(c.get("name",""))}</a>'

        # Pre-build conditional sections (avoid backslashes in f-string expressions for Python 3.11)
        insight_section = f"<p>{client_insight}</p>" if client_insight else ""
        pains_section = (
            "<section class='section section-alt'><div class='container'>"
            "<h2>Виклики вашого бізнесу</h2><div class='grid'>"
            + pains_html + "</div></div></section>"
        ) if pains_html else ""
        mtp_fit_style = 'style="text-align:center;max-width:700px;margin:0 auto 2rem;font-size:1.05rem;color:#555"'
        benefits_section = (
            "<section class='section'><div class='container'>"
            "<h2>Чому MTP Fulfillment</h2>"
            f"<p {mtp_fit_style}>{mtp_fit}</p>"
            "<div class='grid'>" + benefits_html + "</div></div></section>"
        ) if benefits_html else ""
        social_style = 'style="text-align:center;opacity:0.8;margin-bottom:1rem"'
        mtp_years = MTP_COMPANY.get('years', '7+')
        mtp_shipments = MTP_COMPANY.get('shipments_per_month', '60 000+')
        mtp_website = MTP_COMPANY.get('website', '')
        mtp_email = MTP_COMPANY.get('email', '')
        mtp_phone = MTP_COMPANY.get('phone', '')
        mtp_phone_raw = MTP_COMPANY.get('phone_raw', '')

        html = f"""<!DOCTYPE html>
<html lang="uk">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Комерційна Пропозиція для {client_name} від MTP Fulfillment</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap" rel="stylesheet">
<style>
:root {{
  --brand-primary: {brand_primary};
  --brand-secondary: {brand_secondary};
  --mtp-accent: {mtp_accent};
  --hero-text: {hero_text_color};
  --text-dark: #1a1a1a;
  --text-light: #f8f8f8;
  --neutral-bg: #f5f5f5;
}}
* {{ box-sizing: border-box; margin: 0; padding: 0; }}
body {{ font-family: 'Inter', sans-serif; line-height: 1.6; color: var(--text-dark); background: #fff; }}
a {{ color: var(--brand-secondary); text-decoration: none; }}
a:hover {{ text-decoration: underline; }}
h2 {{ color: var(--brand-primary); text-align: center; font-size: 2rem; margin-bottom: 2rem; }}
section:not(.hero):not(.cta-section) h2 {{ color: var(--text-dark) !important; }}
.container {{ max-width: 1100px; margin: 0 auto; padding: 0 1.5rem; }}
.btn {{ display: inline-block; padding: 14px 32px; background: var(--mtp-accent); color: #fff; border: none;
  border-radius: 8px; font-size: 1rem; font-weight: 600; cursor: pointer; transition: transform 0.2s;
  white-space: nowrap; max-width: 100%; text-overflow: ellipsis; }}
.btn:hover {{ transform: translateY(-2px); text-decoration: none; }}

/* Navbar */
.navbar {{ position: fixed; top: 0; left: 0; right: 0; z-index: 100; background: var(--brand-primary);
  padding: 0 2rem; height: 64px; display: flex; align-items: center; justify-content: space-between; }}
.navbar a {{ color: var(--hero-text); font-weight: 600; }}
.navbar .nav-logo {{ font-size: 1.1rem; letter-spacing: 1px; }}
.navbar .nav-actions {{ display: flex; align-items: center; gap: 1.5rem; }}
.navbar .nav-actions a {{ font-size: 0.9rem; opacity: 0.9; }}
.navbar .nav-btn {{ background: var(--mtp-accent); padding: 8px 20px; border-radius: 6px; font-size: 0.85rem; }}

/* Hero */
.hero {{ background: var(--brand-primary); color: var(--hero-text); padding: 120px 2rem 60px; text-align: center; }}
.hero h1 {{ font-size: 2.5rem; font-weight: 700; margin-bottom: 1rem; }}
.hero p {{ font-size: 1.15rem; opacity: 0.9; max-width: 700px; margin: 0 auto 2rem; }}

/* Sections */
.section {{ padding: 60px 0; }}
.section-alt {{ background: var(--neutral-bg); }}
.section-dark {{ background: var(--brand-secondary); color: var(--text-light); }}
.section-dark h2 {{ color: var(--text-light); }}

/* Pain cards */
.grid {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(280px, 1fr)); gap: 1.5rem; }}
.pain-card {{ background: #fff; border-radius: 12px; padding: 1.5rem; box-shadow: 0 2px 12px rgba(0,0,0,0.08);
  border-top: 4px solid var(--mtp-accent); }}
.pain-card h3 {{ font-size: 1.05rem; margin-bottom: 0.5rem; color: var(--text-dark); }}
.pain-card p {{ font-size: 0.95rem; color: #555; }}

/* Benefit cards */
.benefit-card {{ background: #fff; border-radius: 12px; padding: 1.5rem; box-shadow: 0 2px 12px rgba(0,0,0,0.08);
  border-left: 4px solid var(--brand-primary); }}
.benefit-card h3 {{ font-size: 1.05rem; margin-bottom: 0.5rem; color: var(--brand-primary); }}
.benefit-card p {{ font-size: 0.9rem; color: #666; font-style: italic; }}

/* Social proof */
.clients-row {{ display: flex; flex-wrap: wrap; justify-content: center; gap: 2rem; margin-top: 1.5rem; }}
.client-logo {{ background: rgba(255,255,255,0.15); padding: 12px 24px; border-radius: 8px;
  font-weight: 600; font-size: 1rem; color: var(--text-light); }}

/* Tariffs */
.tariff-table {{ width: 100%; border-collapse: collapse; margin-top: 1rem; background: #fff; border-radius: 8px; overflow: hidden; }}
.tariff-table th {{ background: var(--brand-primary); color: var(--hero-text); padding: 12px 16px; text-align: left; font-size: 0.9rem; }}
.tariff-table td {{ padding: 10px 16px; border-bottom: 1px solid #eee; font-size: 0.95rem; }}
.tariff-table td:last-child {{ text-align: right; font-weight: 600; color: var(--brand-primary); }}
.tariff-table tr:nth-child(even) {{ background: var(--neutral-bg); }}
.tariff-table tr.total-row {{ background: #e8f0fe; font-weight: 700; }}

/* CTA */
.cta-section {{ background: var(--brand-primary); padding: 60px 2rem; text-align: center; color: var(--hero-text); }}
.cta-section h2 {{ color: var(--hero-text); }}
.cta-section p {{ font-size: 1.1rem; opacity: 0.9; margin-bottom: 2rem; max-width: 600px; margin-left: auto; margin-right: auto; }}

/* Footer */
.footer {{ background: #0d1b2a; color: #aaa; padding: 40px 2rem; text-align: center; font-size: 0.85rem; }}
.footer a {{ color: #ddd; }}
.footer .footer-contacts {{ margin-bottom: 1rem; font-size: 0.95rem; color: #ccc; }}

/* Mobile */
@media (max-width: 768px) {{
  .hero {{ padding: 100px 1rem 40px; }}
  .hero h1 {{ font-size: 1.6rem; }}
  .hero p {{ font-size: 1rem; }}
  h2 {{ font-size: 1.5rem; }}
  .grid {{ grid-template-columns: 1fr; }}
  .navbar .nav-actions {{ gap: 0.8rem; }}
  .navbar .nav-actions a:not(.nav-btn) {{ display: none; }}
  .section {{ padding: 40px 0; }}
}}
@media (max-width: 480px) {{
  .navbar {{ padding: 0 1rem; }}
  .navbar .nav-btn {{ padding: 6px 14px; font-size: 0.8rem; }}
  .hero h1 {{ font-size: 1.3rem; }}
  .tariff-table td, .tariff-table th {{ padding: 8px 10px; font-size: 0.85rem; }}
}}
</style>
</head>
<body>

<nav class="navbar">
  <a href="{mtp_website}" target="_blank" class="nav-logo">MTP Fulfillment</a>
  <div class="nav-actions">
    <a href="tel:{mtp_phone_raw}">{mtp_phone}</a>
    <a href="{calendly_url}" target="_blank" class="nav-btn btn">Записатись</a>
  </div>
</nav>

<section class="hero">
  <h1>{hook}</h1>
  {insight_section}
  <a href="{calendly_url}" target="_blank" class="btn">Обговорити на Zoom</a>
</section>

{pains_section}

{benefits_section}

<section class="section section-dark">
  <div class="container">
    <h2>Нам довіряють</h2>
    <p {social_style}>{mtp_years} років на ринку &bull; {mtp_shipments} відправок/міс</p>
    <div class="clients-row">{clients_html}</div>
  </div>
</section>

<section class="section section-alt">
  <div class="container">
    <h2>Тарифи</h2>
    <table class="tariff-table">
      <tr><th>Послуга</th><th style="text-align:right">Ціна</th></tr>
      {tariffs_html}
    </table>
    {pricing_html}
  </div>
</section>

<section class="cta-section">
  <h2>Готові масштабуватись?</h2>
  <p>{zoom_cta}</p>
  <a href="{calendly_url}" target="_blank" class="btn" style="font-size:1.1rem;padding:16px 40px">Записатись на Zoom</a>
</section>

<footer class="footer">
  <div class="footer-contacts">
    <a href="mailto:{mtp_email}">{mtp_email}</a> &nbsp;|&nbsp;
    <a href="tel:{mtp_phone_raw}">{mtp_phone}</a> &nbsp;|&nbsp;
    <a href="{mtp_website}" target="_blank">{mtp_website}</a>
  </div>
  &copy; 2026 MTP Fulfillment. Всі права захищені.
</footer>

</body>
</html>"""
        return html

    def _generate_email(self, lead: Lead, analysis: Dict[str, Any], path: str):
        subject = analysis.get("email_subject", f"{lead.name}, пропозиція від MTP Fulfillment")
        opening = analysis.get("email_opening", f"Привіт! Ми подивились на {lead.name} і бачимо великий потенціал.")
        client_insight = analysis.get("client_insight", "")
        mtp_fit = analysis.get("mtp_fit", "")
        zoom_cta = analysis.get("zoom_cta", "Запишіться на безкоштовну Zoom-консультацію!")

        # Build pain points text
        pain_lines = []
        for pain in analysis.get("pain_points", []):
            if isinstance(pain, dict):
                pain_lines.append(f"— {pain.get('title', '')}: {pain.get('description', '')}")
            else:
                pain_lines.append(f"— {pain}")
        pains_text = "\n".join(pain_lines)

        # Build benefits text
        benefit_lines = []
        for kb in analysis.get("key_benefits", []):
            if isinstance(kb, dict):
                benefit_lines.append(f"✓ {kb.get('benefit', '')} — {kb.get('proof', '')}")
            else:
                benefit_lines.append(f"✓ {kb}")
        benefits_text = "\n".join(benefit_lines)

        body = f"""Тема: {subject}

{opening}

{client_insight}

Знайомі ситуації?
{pains_text}

{mtp_fit}

Що ви отримаєте з МТП:
{benefits_text}

{zoom_cta}

У додатку — детальна комерційна пропозиція з тарифами та кошторисом для {lead.name}.

Микола
MTP Fulfillment
mtpgrouppromo@gmail.com | +38 (050) 144-46-45
fulfillmentmtp.com.ua | @nikolay_mtp
"""

        with open(path, "w", encoding="utf-8") as f:
            f.write(body)
